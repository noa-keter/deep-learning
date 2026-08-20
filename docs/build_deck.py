"""
Build the 15-minute project presentation from the frozen final report v3.

Every number on a slide is taken from docs/final_report_v3.docx, which is in turn
backed by the committed JSON under results/figures/. The deck is generated rather
than hand-drawn so it can be rebuilt after any report correction:

    python docs/slide_figures.py     # first, the two slide-only figures
    python docs/build_deck.py

Speaking split, approved 2026-08-20: Ido presents slides 1-8, Noa presents 9-15,
one handover at the method/results boundary. Backup slides follow slide 15 and are
not part of the 13:05 running time.
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

__all__ = ["build_deck", "main"]

# ── Canvas ────────────────────────────────────────────────────────────────────

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.62
CONTENT_W_IN = SLIDE_W_IN - 2 * MARGIN_IN

KICKER_Y_IN = 0.40
TITLE_Y_IN = 0.76
BODY_Y_IN = 1.86

# ── Palette ───────────────────────────────────────────────────────────────────

INK = RGBColor(0x14, 0x1A, 0x2E)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x2E, 0x3A, 0x59)
MUTED = RGBColor(0x6B, 0x76, 0x91)
ACCENT = RGBColor(0xD9, 0x48, 0x4E)
TEAL = RGBColor(0x1C, 0x72, 0x93)
GOLD = RGBColor(0xE8, 0xA3, 0x3D)
CARD = RGBColor(0xF3, 0xF5, 0xF9)
CARD_DARK = RGBColor(0x22, 0x2A, 0x44)
PAPER_DIM = RGBColor(0xC4, 0xCB, 0xDC)

# ── Type ──────────────────────────────────────────────────────────────────────

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

TITLE_PT = 34
KICKER_PT = 12
BODY_PT = 15
SMALL_PT = 12.5
STAT_PT = 46

KICKER_TRACKING = "180"
"""Letter spacing for kicker labels, in hundredths of a point."""

CARD_CORNER = 0.06
"""Rounded-rectangle corner adjustment; PowerPoint's default 0.16 is too soft here."""

REPO_URL = "github.com/noa-keter/deep-learning"

DOCS_DIR = Path(__file__).resolve().parent
ASSETS_DIR = DOCS_DIR / "slide_assets"
FIGURES_DIR = DOCS_DIR.parent / "results" / "figures"
OUTPUT_PATH = DOCS_DIR / "final_presentation.pptx"


# ── Primitives ────────────────────────────────────────────────────────────────


def _blank(prs: Presentation, background: RGBColor):
    """
    Append a blank slide painted with a flat background.

    Args:
        prs: The presentation being built.
        background: Fill color for the whole slide.

    Returns:
        The new slide.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = background
    return slide


def _text(
    slide,
    body: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = BODY_PT,
    color: RGBColor = SLATE,
    font: str = BODY_FONT,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing: float = 1.22,
    space_after_pt: float = 0.0,
    tracking: str | None = None,
):
    """
    Place a text box whose paragraphs are the newline-separated lines of `body`.

    Args:
        slide: Target slide.
        body: Text; each line becomes its own paragraph.
        x: Left edge in inches.
        y: Top edge in inches.
        w: Width in inches.
        h: Height in inches.
        size: Font size in points.
        color: Font color.
        font: Typeface name.
        bold: Whether to bold every run.
        italic: Whether to italicise every run.
        align: Paragraph alignment.
        anchor: Vertical anchor within the box.
        line_spacing: Multiple-of-single line spacing.
        space_after_pt: Trailing space per paragraph, in points.
        tracking: Letter spacing in hundredths of a point, or None.

    Returns:
        The created textbox shape.
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    for index, line in enumerate(body.split("\n")):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        if space_after_pt:
            para.space_after = Pt(space_after_pt)
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
        if tracking is not None:
            run._r.get_or_add_rPr().set("spc", tracking)
    return box


def _card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = CARD,
    outline: RGBColor | None = None,
    outline_pt: float = 1.5,
):
    """
    Draw a rounded rectangle used as a content card.

    Args:
        slide: Target slide.
        x: Left edge in inches.
        y: Top edge in inches.
        w: Width in inches.
        h: Height in inches.
        fill: Interior color.
        outline: Border color, or None for no border.
        outline_pt: Border width in points.

    Returns:
        The created shape.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.adjustments[0] = CARD_CORNER
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        shape.line.width = Pt(outline_pt)
    shape.shadow.inherit = False
    if shape.has_text_frame:
        shape.text_frame.text = ""
    return shape


def _chrome(
    slide,
    kicker: str,
    title: str,
    *,
    dark: bool = False,
    title_color: RGBColor | None = None,
    title_pt: float = TITLE_PT,
) -> None:
    """
    Lay down the repeating slide furniture: kicker label and title.

    The uppercase tracked kicker is the deck's one visual motif; it also tells a
    grader which of the seven required presentation elements the slide serves.

    Args:
        slide: Target slide.
        kicker: Short uppercase section label.
        title: Slide title, newline-separated for multi-line titles.
        dark: Whether the slide background is dark.
        title_color: Overrides the default title color.
        title_pt: Title font size in points.
    """
    _text(
        slide,
        kicker.upper(),
        MARGIN_IN,
        KICKER_Y_IN,
        CONTENT_W_IN,
        0.3,
        size=KICKER_PT,
        color=ACCENT,
        bold=True,
        tracking=KICKER_TRACKING,
    )
    _text(
        slide,
        title,
        MARGIN_IN,
        TITLE_Y_IN,
        CONTENT_W_IN,
        0.98,
        size=title_pt,
        color=title_color or (PAPER if dark else INK),
        font=HEAD_FONT,
        bold=True,
        line_spacing=1.08,
    )


def _stat(
    slide,
    value: str,
    label: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    value_color: RGBColor = INK,
    label_color: RGBColor = MUTED,
    fill: RGBColor | None = CARD,
    value_pt: float = STAT_PT,
) -> None:
    """
    Draw a large-number callout with a caption underneath.

    Args:
        slide: Target slide.
        value: The headline figure, e.g. "0.977".
        label: Caption below it; newline-separated for multiple lines.
        x: Left edge in inches.
        y: Top edge in inches.
        w: Width in inches.
        h: Height in inches.
        value_color: Color of the big number.
        label_color: Color of the caption.
        fill: Card fill, or None to draw the stat with no card behind it.
        value_pt: Font size of the big number, in points.
    """
    if fill is not None:
        _card(slide, x, y, w, h, fill=fill)
    _text(
        slide,
        value,
        x + 0.24,
        y + 0.20,
        w - 0.48,
        h * 0.52,
        size=value_pt,
        color=value_color,
        font=HEAD_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    _text(
        slide,
        label,
        x + 0.20,
        y + h * 0.55,
        w - 0.40,
        h * 0.40,
        size=SMALL_PT,
        color=label_color,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.TOP,
        line_spacing=1.16,
    )


def _bullets(
    slide,
    items: list[tuple[str, str]],
    x: float,
    y: float,
    w: float,
    *,
    row_h: float = 0.86,
    gap: float = 0.14,
    dark: bool = False,
    marker_color: RGBColor = ACCENT,
    lead_color: RGBColor | None = None,
) -> float:
    """
    Draw rows of "bold lead — supporting sentence", each with a small filled marker.

    Args:
        slide: Target slide.
        items: (lead, detail) pairs; detail may be empty.
        x: Left edge in inches.
        y: Top edge of the first row, in inches.
        w: Row width in inches.
        row_h: Height of each row in inches.
        gap: Vertical gap between rows in inches.
        dark: Whether the slide background is dark.
        marker_color: Fill color of the square marker.
        lead_color: Color of the bold lead text; defaults by background.

    Returns:
        The y coordinate just below the last row, in inches.
    """
    marker_side = 0.13
    text_x = x + 0.42
    cursor = y
    for lead, detail in items:
        marker = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(cursor + 0.10),
            Inches(marker_side),
            Inches(marker_side),
        )
        marker.adjustments[0] = 0.3
        marker.fill.solid()
        marker.fill.fore_color.rgb = marker_color
        marker.line.fill.background()
        marker.shadow.inherit = False

        _text(
            slide,
            lead,
            text_x,
            cursor,
            w - 0.42,
            0.34,
            size=BODY_PT + 1,
            color=lead_color or (PAPER if dark else INK),
            bold=True,
        )
        if detail:
            _text(
                slide,
                detail,
                text_x,
                cursor + 0.36,
                w - 0.42,
                row_h - 0.36,
                size=BODY_PT - 1,
                color=PAPER_DIM if dark else SLATE,
                line_spacing=1.20,
            )
        cursor += row_h + gap
    return cursor


def _grid_table(
    slide,
    header: list[str],
    rows: list[list[str]],
    x: float,
    y: float,
    w: float,
    *,
    col_ratios: list[float] | None = None,
    row_h: float = 0.44,
    header_h: float = 0.44,
    dark: bool = False,
    emphasis_col: int | None = None,
    emphasis_color: RGBColor = ACCENT,
    body_pt: float = BODY_PT - 1,
    header_pt: float = SMALL_PT,
) -> float:
    """
    Draw a lightweight table from text boxes and banding rectangles.

    python-pptx's native table carries a heavy default style that has to be
    stripped cell by cell, so the deck draws its own instead.

    Args:
        slide: Target slide.
        header: Header cell texts.
        rows: Body rows, each the same length as `header`.
        x: Left edge in inches.
        y: Top edge in inches.
        w: Total width in inches.
        col_ratios: Relative column widths; defaults to equal.
        row_h: Body row height in inches.
        header_h: Header row height in inches.
        dark: Whether the slide background is dark.
        emphasis_col: Index of a column to colour, or None.
        emphasis_color: Colour applied to `emphasis_col`.
        body_pt: Body cell font size in points.
        header_pt: Header cell font size in points.

    Returns:
        The y coordinate just below the table, in inches.

    Raises:
        ValueError: If a row's length does not match the header's.
    """
    n_cols = len(header)
    for row in rows:
        if len(row) != n_cols:
            raise ValueError(f"row {row!r} has {len(row)} cells, expected {n_cols}")

    ratios = col_ratios or [1.0] * n_cols
    total = sum(ratios)
    widths = [w * r / total for r in ratios]
    lefts = [x + sum(widths[:i]) for i in range(n_cols)]

    body_color = PAPER_DIM if dark else SLATE
    head_color = MUTED if dark else MUTED
    band = CARD_DARK if dark else CARD

    for i, cell in enumerate(header):
        _text(
            slide,
            cell,
            lefts[i] + 0.14,
            y + 0.08,
            widths[i] - 0.28,
            header_h - 0.10,
            size=header_pt,
            color=head_color,
            bold=True,
            align=PP_ALIGN.CENTER if i else PP_ALIGN.LEFT,
            tracking="60",
        )

    cursor = y + header_h
    for r_index, row in enumerate(rows):
        if r_index % 2 == 0:
            _card(slide, x, cursor, w, row_h, fill=band)
        for i, cell in enumerate(row):
            is_emphasis = emphasis_col is not None and i == emphasis_col
            _text(
                slide,
                cell,
                lefts[i] + 0.14,
                cursor,
                widths[i] - 0.28,
                row_h,
                size=body_pt,
                color=emphasis_color if is_emphasis else body_color,
                bold=(i == 0) or is_emphasis,
                font=BODY_FONT,
                align=PP_ALIGN.CENTER if i else PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE,
            )
        cursor += row_h
    return cursor


def _picture(
    slide, path: Path, x: float, y: float, w: float, *, max_h: float | None = None
) -> float:
    """
    Place an image at a given width, shrinking and re-centring it if it would be too tall.

    The height cap is what keeps a wide figure from running off the bottom of the
    slide; without it the only symptom is a caption silently cut in half.

    Args:
        slide: Target slide.
        path: Image file.
        x: Left edge in inches.
        y: Top edge in inches.
        w: Desired width in inches.
        max_h: Largest permitted rendered height in inches, or None for no cap.

    Returns:
        The rendered height in inches.

    Raises:
        FileNotFoundError: If the image is missing.
    """
    if not path.exists():
        raise FileNotFoundError(f"slide asset not found: {path}")

    probe = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    height = Emu(probe.height).inches
    if max_h is not None and height > max_h:
        scale = max_h / height
        width = w * scale
        probe.width = Inches(width)
        probe.height = Inches(max_h)
        probe.left = Inches(x + (w - width) / 2.0)
        height = max_h
    return height


def _notes(slide, cues: str) -> None:
    """
    Attach short delivery cues to a slide.

    Args:
        slide: Target slide.
        cues: Plain-text cue lines.
    """
    slide.notes_slide.notes_text_frame.text = cues


# ── Slides ────────────────────────────────────────────────────────────────────


def _slide_01_title(prs: Presentation) -> None:
    """
    Title slide: project name, both presenters, course, repository.
    """
    slide = _blank(prs, INK)
    _text(
        slide,
        "DEEP LEARNING  ·  TEL AVIV UNIVERSITY  ·  2026B",
        MARGIN_IN,
        1.42,
        CONTENT_W_IN,
        0.3,
        size=KICKER_PT,
        color=ACCENT,
        bold=True,
        tracking=KICKER_TRACKING,
    )
    _text(
        slide,
        "Resolution Bias and\nCross-Generator Detection\nof AI-Generated Images",
        MARGIN_IN,
        2.00,
        CONTENT_W_IN - 1.0,
        2.9,
        size=42,
        color=PAPER,
        font=HEAD_FONT,
        bold=True,
        line_spacing=1.14,
    )
    _text(
        slide,
        "Noa Keter  ·  Ido Josephsberg",
        MARGIN_IN,
        5.50,
        6.6,
        0.34,
        size=17,
        color=PAPER,
        bold=True,
    )
    _text(
        slide,
        "Lecturer: Dr. Shay Maymon",
        MARGIN_IN,
        5.94,
        6.6,
        0.30,
        size=13.5,
        color=PAPER_DIM,
    )
    _text(
        slide,
        REPO_URL,
        MARGIN_IN,
        6.34,
        6.6,
        0.30,
        size=13,
        color=PAPER_DIM,
        font=MONO_FONT,
    )
    _stat(
        slide,
        "56",
        "training runs\n4 strategies x 7 generators x 2 seeds",
        9.30,
        4.62,
        3.40,
        1.94,
        value_color=GOLD,
        label_color=PAPER_DIM,
        fill=CARD_DARK,
        value_pt=42,
    )
    _notes(
        slide,
        "0:15  |  Ido\n"
        "- Title, both names, repo is public and linked.\n"
        "- One line only: detecting AI-generated images, and what the benchmark "
        "is really measuring.\n"
        "- Do not explain the 56 here; it is a hook for slide 8.",
    )


def _slide_02_problem(prs: Presentation) -> None:
    """
    Motivation 1: the in-domain / cross-generator gap.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Motivation  ·  1 of 3",
        "A detector that looks excellent, until the generator changes",
        title_pt=31,
    )
    _text(
        slide,
        "Train a CNN to separate real photographs from images made by one generator, and it "
        "scores almost perfectly on that generator. Point it at a generator it has never seen "
        "and the accuracy falls off a cliff.",
        MARGIN_IN,
        BODY_Y_IN,
        7.05,
        1.15,
        size=BODY_PT + 0.5,
        color=SLATE,
        line_spacing=1.32,
    )
    _bullets(
        slide,
        [
            (
                "The question this project asks",
                "How much of that in-distribution success reflects the generative process,\n"
                "and how much reflects the way the benchmark was built?",
            ),
            (
                "This gap is the published pattern, not our finding",
                "What is ours is the explanation — and it is not about the generators.",
            ),
        ],
        MARGIN_IN,
        3.42,
        7.05,
        row_h=1.14,
        gap=0.46,
    )
    _stat(slide, "0.908", "same generator\nin-domain accuracy", 8.10, 2.22, 2.28, 1.78, value_color=TEAL)
    _stat(slide, "0.655", "unseen generator\ncross-generator accuracy", 10.72, 2.22, 2.28, 1.78, value_color=ACCENT)
    _text(
        slide,
        "↓",
        8.10,
        4.14,
        4.90,
        0.56,
        size=26,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _card(slide, 8.10, 4.78, 4.90, 1.00, fill=CARD)
    _text(
        slide,
        "a 25-point drop, from one change:\nthe generator at test time",
        8.30,
        5.00,
        4.50,
        0.64,
        size=13.5,
        color=SLATE,
        align=PP_ALIGN.CENTER,
        line_spacing=1.24,
    )
    _text(
        slide,
        "center_crop arm, mean of two seeds, n = 1,000 balanced images per cell",
        8.10,
        5.96,
        4.90,
        0.3,
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "0:45  |  Ido\n"
        "- Set the scene: this gap is the known, published pattern, not our finding.\n"
        "- 0.908 -> 0.655 is our own center_crop arm, so the audience meets our numbers early.\n"
        "- End on the question verbatim: process, or benchmark construction?\n"
        "- Do not mention size yet. That is the next slide's reveal.",
    )


def _slide_03_confound(prs: Presentation) -> None:
    """
    Motivation 2: the size confound and the zero-parameter rule that exploits it.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Motivation  ·  2 of 3",
        "Every generator emits one fixed size. Real photographs do not.",
        title_pt=29,
    )
    _grid_table(
        slide,
        ["", "BigGAN", "ADM  GLIDE  VQDM", "SD 1.5  Wukong", "Midjourney", "Real"],
        [["Native resolution", "128²", "256²", "512²", "1024²", "variable"]],
        MARGIN_IN,
        BODY_Y_IN,
        CONTENT_W_IN,
        col_ratios=[1.55, 1.0, 1.85, 1.45, 1.15, 1.15],
        row_h=0.52,
        header_h=0.42,
    )
    _text(
        slide,
        "So the dimensions alone identify both the class and the generator. A detector can score "
        "well without ever examining pixel content.",
        MARGIN_IN,
        3.10,
        7.15,
        0.86,
        size=BODY_PT + 0.5,
        color=SLATE,
        line_spacing=1.30,
    )
    _card(slide, MARGIN_IN, 4.14, 7.15, 1.90, fill=CARD)
    _text(
        slide,
        "The shortcut, made explicit",
        MARGIN_IN + 0.30,
        4.36,
        6.55,
        0.32,
        size=13,
        color=MUTED,
        bold=True,
        tracking="60",
    )
    _text(
        slide,
        "“Predict synthetic whenever the image is square.”\n"
        "Zero parameters. No training. No pixels.",
        MARGIN_IN + 0.30,
        4.76,
        6.55,
        1.00,
        size=BODY_PT + 1,
        color=INK,
        line_spacing=1.34,
    )
    _stat(
        slide,
        "0.977",
        "accuracy on every one of the 49\nsource → target pairs",
        8.28,
        3.16,
        4.42,
        2.00,
        value_color=ACCENT,
        value_pt=54,
    )
    _text(
        slide,
        "Higher than the cross-generator accuracy\nof every trained detector in this talk.",
        8.28,
        5.36,
        4.42,
        0.72,
        size=14,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
        line_spacing=1.26,
    )
    _notes(
        slide,
        "1:05  |  Ido\n"
        "- Walk the table left to right, then land the punchline: size identifies class AND generator.\n"
        "- Read the rule aloud. Emphasise zero parameters, no training, no pixels.\n"
        "- 0.977 beats every trained detector we will show. Let that sit for a beat.\n"
        "- This is the slide the whole talk hangs on. Do not rush it.",
    )


def _slide_04_question(prs: Presentation) -> None:
    """
    Motivation 3: the framing that makes the correction the experimental variable.
    """
    slide = _blank(prs, INK)
    _chrome(
        slide,
        "Motivation  ·  3 of 3",
        "Removing the cue is not a neutral act",
        dark=True,
    )
    _bullets(
        slide,
        [
            (
                "The fingerprint is high-frequency",
                "Cropping discards context. Rescaling attenuates the very frequencies that\n"
                "carry the signal. Padding introduces a border of its own.",
            ),
            (
                "Prior work applies one correction and reports the corrected number",
                "It does not ask whether the choice of correction changes the answer.",
            ),
            (
                "We make the correction itself the experimental variable",
                "One architecture, one dataset, one protocol, four ways of equalizing size.",
            ),
        ],
        MARGIN_IN,
        BODY_Y_IN + 0.16,
        8.05,
        row_h=1.02,
        gap=0.60,
        dark=True,
    )
    _card(slide, 9.10, 3.14, 3.62, 1.86, fill=CARD_DARK)
    _text(
        slide,
        "Do the conclusions\nsurvive that choice?",
        9.34,
        3.40,
        3.14,
        0.96,
        size=18,
        color=PAPER,
        font=HEAD_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
        line_spacing=1.22,
    )
    _text(
        slide,
        "They do not.",
        9.34,
        4.42,
        3.14,
        0.44,
        size=22,
        color=ACCENT,
        font=HEAD_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "0:35  |  Ido\n"
        "- Three beats: the fingerprint is fragile, prior work corrects once, we vary the correction.\n"
        "- Deliver 'They do not.' flatly and move on. It is the thesis; slide 11 is the evidence.\n"
        "- Keep this short, it is a hinge slide.",
    )


def _slide_05_related(prs: Presentation) -> None:
    """
    Required element: brief review of related work.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Related work",
        "Three papers, and where we start",
    )
    entries = [
        (
            "2020",
            "Wang et al.  —  CNN-generated images are surprisingly easy to spot",
            "Established the pattern: easy in-distribution, uneven transfer to unseen architectures.",
            TEAL,
        ),
        (
            "2023",
            "GenImage  —  a million-scale benchmark",
            "Scaled the question and formalized the cross-generator protocol we adopt: "
            "train on one generator, evaluate on all.",
            TEAL,
        ),
        (
            "2024",
            "Grommelt et al.  —  Fake or JPEG?",
            "Showed GenImage carries two construction artifacts, JPEG compression and image size. "
            "Removing both adds >11 points cross-generator.",
            ACCENT,
        ),
    ]
    cursor = BODY_Y_IN
    for year, title, detail, color in entries:
        _card(slide, MARGIN_IN, cursor, 1.02, 0.62, fill=color)
        _text(
            slide,
            year,
            MARGIN_IN,
            cursor,
            1.02,
            0.62,
            size=15,
            color=PAPER,
            font=HEAD_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _text(slide, title, MARGIN_IN + 1.28, cursor + 0.02, 10.35, 0.34, size=BODY_PT + 1, color=INK, bold=True)
        _text(
            slide,
            detail,
            MARGIN_IN + 1.28,
            cursor + 0.40,
            10.35,
            0.62,
            size=BODY_PT - 1,
            color=SLATE,
            line_spacing=1.22,
        )
        cursor += 1.22

    _card(slide, MARGIN_IN, 5.62, CONTENT_W_IN, 1.06, fill=CARD)
    _text(
        slide,
        "We take the size bias as established and begin where that work stops. It corrects the bias "
        "once, by a single mechanism. Our claim is different: the correction is itself a free "
        "parameter, and it reorders published results.",
        MARGIN_IN + 0.30,
        5.82,
        CONTENT_W_IN - 0.60,
        0.72,
        size=14,
        color=INK,
        line_spacing=1.26,
    )
    _notes(
        slide,
        "0:45  |  Ido\n"
        "- One sentence per paper, do not linger. The audience needs the chain, not the detail.\n"
        "- Grommelt is the one that matters: it proves the bias is real and worth 11 points.\n"
        "- Land on the bottom card: they correct once, we vary the correction.",
    )


def _slide_06_data(prs: Presentation) -> None:
    """
    Required element: description of the data.
    """
    slide = _blank(prs, PAPER)
    _chrome(slide, "Method  ·  Data", "Tiny-GenImage")
    stats = [
        ("35,000", "images\n8.4 GB", INK),
        ("7", "generators\n2,500 images each", INK),
        ("17,500", "real ImageNet\nphotographs", INK),
        ("1 of 7", "generators is PNG\nthe rest are JPEG", TEAL),
    ]
    for index, (value, label, color) in enumerate(stats):
        _stat(
            slide,
            value,
            label,
            MARGIN_IN + index * 3.10,
            BODY_Y_IN,
            2.86,
            1.66,
            value_color=color,
            value_pt=34,
        )
    _bullets(
        slide,
        [
            (
                "Format bias is largely absent here, which is what makes the subset usable",
                "GenImage's JPEG/PNG confound is gone; ADM alone is PNG. Size is the only "
                "substantial surviving confound — exactly the one we want to study.",
            ),
            (
                "Splits, and why no reported number ever influenced training",
                "3,600 train / 400 validation per run, from the shipped train split. The shipped "
                "validation split is held out as our test set: 4,000 images. Model selection uses "
                "the internal validation set only.",
            ),
            (
                "Every matrix cell is exactly balanced",
                "500 target fakes + the same fixed 500 real images, so n = 1,000 and accuracy is "
                "directly interpretable. Binomial standard error ±1.6 pp.",
            ),
        ],
        MARGIN_IN,
        3.86,
        CONTENT_W_IN,
        row_h=0.90,
        gap=0.10,
    )
    _notes(
        slide,
        "0:45  |  Ido\n"
        "- Four tiles fast. Only the fourth needs a sentence: the format confound is gone.\n"
        "- Stress that the test split never touches training or model selection.\n"
        "- Balanced cells matter because every number later is a plain accuracy.\n"
        "- Licence is CC BY-NC-SA 4.0 if anyone asks.",
    )


def _slide_07_strategies(prs: Presentation) -> None:
    """
    Required element: the method, part one. The four size-equalization strategies.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Method  ·  Preprocessing",
        "Four ways to make every image 128 × 128",
    )
    _text(
        slide,
        "The first two preserve native pixel statistics and discard content. The last two resample "
        "— and pad adds a border.",
        MARGIN_IN,
        1.74,
        CONTENT_W_IN,
        0.34,
        size=14,
        color=SLATE,
    )
    _picture(
        slide,
        ASSETS_DIR / "strategies_diagram.png",
        MARGIN_IN,
        2.14,
        CONTENT_W_IN,
        max_h=4.90,
    )
    _notes(
        slide,
        "0:55  |  Ido\n"
        "- Top row a generated image, bottom row a real photo. Same four operations.\n"
        "- Crops zoom in: fewer, larger grid cells. Rescale keeps the whole frame and squashes it.\n"
        "- Two red callouts are the payload: pad == rescale on every generator, and only real "
        "photos get bars.\n"
        "- Plant it now, do not explain the consequence. Slide 12 collects it.",
    )


def _slide_08_model(prs: Presentation) -> None:
    """
    Required elements: the model and hyperparameters, and the experimental design.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Method  ·  Model + Experiments",
        "One architecture, held fixed, 56 times",
    )
    _bullets(
        slide,
        [
            (
                "Compact CNN, 1.17 M parameters, trained from scratch in PyTorch",
                "Four blocks of [Conv 3×3, BatchNorm, ReLU] ×2 then 2×2 max-pool, "
                "widths 32/64/128/256, global average pooling, dropout 0.3, one linear unit.",
            ),
            (
                "From scratch is required, not stylistic",
                "The real class is ImageNet. An ImageNet-pretrained backbone would arrive having "
                "already seen the negative class.",
            ),
            (
                "Two choices follow from the nature of the signal",
                "Nothing downsamples before the end of block one, because the cue is "
                "high-frequency. The head is global average pooling, because the evidence is a "
                "location-free texture statistic.",
            ),
        ],
        MARGIN_IN,
        BODY_Y_IN,
        7.30,
        row_h=1.06,
        gap=0.16,
    )
    _grid_table(
        slide,
        ["Setting", "Value"],
        [
            ["Optimizer", "AdamW"],
            ["Learning rate", "3×10⁻⁴"],
            ["Weight decay", "1×10⁻⁴"],
            ["Schedule", "cosine to 0,  40 epochs"],
            ["Batch / precision", "128,  AMP float16"],
            ["Augmentation", "horizontal flip, p = 0.5"],
            ["Selection", "best internal validation"],
        ],
        8.22,
        BODY_Y_IN,
        4.49,
        col_ratios=[1.0, 1.42],
        row_h=0.39,
        header_h=0.34,
        body_pt=13,
    )
    _card(slide, 8.22, 5.06, 4.49, 1.28, fill=INK)
    _text(
        slide,
        "4 strategies  ×  7 sources  ×  2 seeds",
        8.44,
        5.24,
        4.05,
        0.30,
        size=13,
        color=PAPER_DIM,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        "56 runs   ·   ~2 GPU-hours",
        8.44,
        5.60,
        4.05,
        0.44,
        size=20,
        color=PAPER,
        font=HEAD_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        "free-tier Colab T4,  127 s per run,  each evaluated on all seven targets",
        8.22,
        6.46,
        4.49,
        0.5,
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        line_spacing=1.2,
    )
    _notes(
        slide,
        "0:55  |  Ido\n"
        "- Architecture in one breath, then spend the time on WHY from scratch: real == ImageNet.\n"
        "- That single fact rules out every pretrained backbone. Expect a question on it.\n"
        "- Config table: point, do not read.\n"
        "- Close on 56 runs / 2 GPU-hours, then hand over to Noa.",
    )


def _slide_09_baselines(prs: Presentation) -> None:
    """
    Results 1, and the handover: the two zero-parameter reference rules.
    """
    slide = _blank(prs, INK)
    _chrome(
        slide,
        "Results  ·  1 of 4",
        "What every trained detector is competing against",
        dark=True,
    )
    _stat(
        slide,
        "0.977",
        "square rule\nthe same score on all 49 cells,\ndiagonal and off-diagonal alike",
        MARGIN_IN,
        2.06,
        3.88,
        2.26,
        value_color=ACCENT,
        label_color=PAPER_DIM,
        fill=CARD_DARK,
        value_pt=48,
    )
    _stat(
        slide,
        "1.000 / 0.619",
        "size-lookup rule\ndiagonal / off-diagonal",
        4.72,
        2.06,
        3.88,
        2.26,
        value_color=GOLD,
        label_color=PAPER_DIM,
        fill=CARD_DARK,
        value_pt=32,
    )
    _card(slide, 8.84, 2.06, 3.87, 2.26, fill=CARD_DARK)
    _text(
        slide,
        "Neither rule measures\nanything about the\ngenerative process.",
        9.08,
        2.42,
        3.39,
        1.10,
        size=17,
        color=PAPER,
        font=HEAD_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
        line_spacing=1.24,
    )
    _text(
        slide,
        "This is the shortcut every\ntrained detector competes against.",
        9.08,
        3.56,
        3.39,
        0.60,
        size=12.5,
        color=PAPER_DIM,
        align=PP_ALIGN.CENTER,
        line_spacing=1.22,
    )
    _bullets(
        slide,
        [
            (
                "Why the square rule works: TPR 1.000, TNR 0.954",
                "Every generated image in this dataset is square, and 95.4 % of real "
                "photographs are not.",
            ),
            (
                "The size-lookup rule exposes the mechanism exactly",
                "It scores 1.000 whenever source and target share a native resolution and exactly "
                "0.500 whenever they do not — in 47 of 49 cells. The two exceptions trace to a "
                "single 512² Midjourney row.",
            ),
        ],
        MARGIN_IN,
        4.66,
        CONTENT_W_IN,
        row_h=0.92,
        gap=0.10,
        dark=True,
    )
    _notes(
        slide,
        "0:40  |  NOA STARTS HERE  — handover\n"
        "- Open by re-anchoring: these two rules have no parameters and no training.\n"
        "- 0.977 flat across all 49 cells; the rule cannot tell in-domain from out-of-domain "
        "because it never looked at a generator.\n"
        "- Size-lookup 1.000 / 0.500 by shared resolution is the cleanest proof the cue is "
        "resolution and nothing else.\n"
        "- Then: so what do the trained detectors do?",
    )


def _slide_10_matrices(prs: Presentation) -> None:
    """
    Results 2: the four transfer matrices.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Results  ·  2 of 4",
        "Four corrections, four different pictures",
    )
    _text(
        slide,
        "Rows are the generator a detector trained on, columns the generator it was evaluated on. "
        "Two seeds, n = 1,000 per cell, one shared colour scale.",
        MARGIN_IN,
        1.72,
        CONTENT_W_IN,
        0.34,
        size=13.5,
        color=SLATE,
    )
    _picture(slide, ASSETS_DIR / "transfer_matrices_wide.png", MARGIN_IN, 2.18, CONTENT_W_IN)
    _card(slide, MARGIN_IN, 5.92, CONTENT_W_IN, 0.92, fill=CARD)
    _text(
        slide,
        "The diagonals are broadly similar — 90.8, 89.2, 90.2. The off-diagonal structure is not, "
        "and that difference is the result. ADM's row is the clearest case: strong under both crops, "
        "at chance under rescale.",
        MARGIN_IN + 0.30,
        6.16,
        CONTENT_W_IN - 0.60,
        0.52,
        size=14,
        color=INK,
        line_spacing=1.24,
    )
    _notes(
        slide,
        "1:05  |  Noa\n"
        "- Orient the audience first: rows train, columns test, diagonal is in-domain.\n"
        "- Point at the three similar diagonals, then at the off-diagonals that are not similar.\n"
        "- Trace ADM's row across the four panels with a finger. 86/87 -> 50s under rescale.\n"
        "- Do not read cells aloud. The pattern is the message.\n"
        "- Pad looks best. Flag that you will come back to it, then move on.",
    )


def _slide_11_headline(prs: Presentation) -> None:
    """
    Results 3: the main result, the instability of the generator ranking.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Results  ·  3 of 4  ·  the main result",
        "The ranking does not survive the correction",
    )
    _picture(slide, FIGURES_DIR / "ranking_report.png", MARGIN_IN, 1.80, 8.16)
    _stat(
        slide,
        "ρ = +0.964",
        "the two non-resampling arms agree\np = 0.003",
        9.02,
        1.86,
        3.70,
        1.44,
        value_color=TEAL,
        value_pt=27,
    )
    _stat(
        slide,
        "ρ = 0.000",
        "rescale retains none of\nthat structure",
        9.02,
        3.44,
        3.70,
        1.44,
        value_color=ACCENT,
        value_pt=27,
    )
    _card(slide, 9.02, 5.02, 3.70, 1.62, fill=INK)
    _text(
        slide,
        "ADM",
        9.24,
        5.22,
        3.26,
        0.34,
        size=14,
        color=PAPER_DIM,
        bold=True,
        align=PP_ALIGN.CENTER,
        tracking="60",
    )
    _text(
        slide,
        "1st  →  7th",
        9.24,
        5.58,
        3.26,
        0.50,
        size=25,
        color=PAPER,
        font=HEAD_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        "best source under both crops,\nworst under rescaling",
        9.24,
        6.12,
        3.26,
        0.44,
        size=11.5,
        color=PAPER_DIM,
        align=PP_ALIGN.CENTER,
        line_spacing=1.2,
    )
    _text(
        slide,
        "Two papers using the same data, architecture and protocol would disagree about which "
        "generator is hardest — purely because one rescaled and the other cropped.",
        MARGIN_IN,
        6.28,
        8.16,
        0.62,
        size=14,
        color=INK,
        bold=True,
        line_spacing=1.24,
    )
    _notes(
        slide,
        "1:25  |  Noa  — the slide that matters most\n"
        "- Left panel: each line is one generator; vertical order IS the ranking.\n"
        "- The crops agree almost perfectly. That agreement is the anchor: the ranking is a real, "
        "reproducible property of the data.\n"
        "- Measured against it, rescale produces an unrelated order. ADM goes first to last.\n"
        "- Honesty, say it out loud: at n = 7, rho = 0.000 is absence of evidence for association, "
        "NOT evidence of independence. The claim rests on the +0.964 and on the size of the move.\n"
        "- Seed disagreement averages 0.012-0.024 and never exceeds 0.076, so a seven-position "
        "inversion is far outside noise.\n"
        "- Close on the bottom line and pause.",
    )


def _slide_12_pad(prs: Presentation) -> None:
    """
    Results 4: the pad arm scores best and is contaminated.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Results  ·  4 of 4",
        "pad scores highest — and that is the failure",
        title_color=ACCENT,
    )
    chain = [
        ("Every generator is square", "so padding changes nothing for them"),
        ("Only real photos get bars", "variable aspect ratio, so they alone are padded"),
        ("The border becomes a perfect proxy", "border ⇔ real, learnable in a few epochs"),
    ]
    for index, (lead, detail) in enumerate(chain):
        top = BODY_Y_IN + index * 1.26
        _card(slide, MARGIN_IN, top, 7.05, 0.98, fill=CARD)
        _text(slide, lead, MARGIN_IN + 0.28, top + 0.14, 6.5, 0.32, size=BODY_PT + 1, color=INK, bold=True)
        _text(slide, detail, MARGIN_IN + 0.28, top + 0.52, 6.5, 0.34, size=BODY_PT - 1.5, color=SLATE)
        if index < len(chain) - 1:
            # The connector needs a box taller than its glyph, or PowerPoint clips it.
            _text(
                slide,
                "↓",
                MARGIN_IN + 0.20,
                top + 0.98,
                0.60,
                0.28,
                size=14,
                color=MUTED,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
    _stat(
        slide,
        "0.977",
        "where a pure border-reader\nwould ceiling:  (1.000 + 0.954) / 2",
        8.10,
        BODY_Y_IN,
        4.60,
        1.62,
        value_color=MUTED,
        value_pt=40,
    )
    _stat(
        slide,
        "0.978",
        "pad's measured diagonal",
        8.10,
        3.62,
        4.60,
        1.42,
        value_color=ACCENT,
        value_pt=40,
    )
    _card(slide, 8.10, 5.22, 4.60, 1.44, fill=INK)
    _text(
        slide,
        "Corroboration",
        8.34,
        5.40,
        4.12,
        0.28,
        size=12,
        color=PAPER_DIM,
        bold=True,
        tracking="60",
    )
    _text(
        slide,
        "Sharing a native resolution helps transfer by\n"
        "+0.113, +0.113, +0.135 in the other three arms\n"
        "— but only +0.007 for pad.",
        8.34,
        5.72,
        4.12,
        0.80,
        size=12,
        color=PAPER,
        line_spacing=1.24,
    )
    _text(
        slide,
        "Higher accuracy here marks a worse experiment, not a better detector. "
        "Padding does not remove the size confound — it re-encodes it as a visible border.",
        MARGIN_IN,
        5.76,
        7.05,
        0.86,
        size=14,
        color=ACCENT,
        bold=True,
        line_spacing=1.28,
    )
    _notes(
        slide,
        "0:55  |  Noa\n"
        "- Walk the three-step chain slowly; it is a proof, not a list.\n"
        "- Then the two numbers side by side: predicted ceiling 0.977, measured 0.978. "
        "That near-identity is the evidence.\n"
        "- Corroboration: pad is the ONLY arm whose transfer is indifferent to generator "
        "similarity, exactly as a border-reader would be.\n"
        "- Say the closing sentence verbatim. It is the project's thesis in one line.\n"
        "- Never call pad the best arm.",
    )


def _slide_13_attribution(prs: Presentation) -> None:
    """
    Analysis 1: input-gradient attribution confirms the border mechanism.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Analysis  ·  1 of 2",
        "The model really is reading its border",
    )
    _text(
        slide,
        "Input-gradient saliency, 200 images per tag, seed 0. Border mass is the fraction of the "
        "mean saliency map falling in the outer 16-pixel ring — against a uniform null of 0.4375.",
        MARGIN_IN,
        1.72,
        CONTENT_W_IN,
        0.56,
        size=13.5,
        color=SLATE,
        line_spacing=1.24,
    )
    _picture(slide, FIGURES_DIR / "attribution_border_mass.png", MARGIN_IN, 2.36, 7.16)
    _grid_table(
        slide,
        ["Strategy", "Real", "Generators", "Real 1st"],
        [
            ["center_crop", "0.4345", "0.4243", "0 of 7"],
            ["random_crop", "0.4310", "0.4175", "1 of 7"],
            ["rescale", "0.3793", "0.3721", "2 of 7"],
            ["pad", "0.5613", "0.4113", "7 of 7"],
        ],
        7.98,
        2.32,
        4.73,
        col_ratios=[1.42, 1.0, 1.30, 1.0],
        row_h=0.44,
        header_h=0.40,
        body_pt=13,
        header_pt=11.5,
    )
    _card(slide, 7.98, 4.50, 4.73, 1.94, fill=INK)
    _text(
        slide,
        "The two groups do not overlap",
        8.22,
        4.70,
        4.25,
        0.32,
        size=13,
        color=PAPER,
        bold=True,
    )
    _text(
        slide,
        "Every pad run exceeds +0.11 over its own generator mean. "
        "No run in any other arm reaches +0.04.\n\n"
        "Padding does not raise edge saliency in general — only the class that "
        "receives a border is affected.",
        8.22,
        5.08,
        4.25,
        1.24,
        size=11.5,
        color=PAPER_DIM,
        line_spacing=1.20,
    )
    _text(
        slide,
        "This confirms from inside the model what the accuracies alone could only imply.",
        MARGIN_IN,
        5.16,
        7.16,
        0.62,
        size=14,
        color=INK,
        bold=True,
        line_spacing=1.24,
    )
    _text(
        slide,
        "rescale also draws on the most centrally concentrated evidence of the four arms, "
        "0.373 against 0.419 to 0.430.",
        MARGIN_IN,
        5.92,
        7.16,
        0.62,
        size=12.5,
        color=MUTED,
        line_spacing=1.22,
    )
    _notes(
        slide,
        "0:50  |  Noa\n"
        "- This is the cue-level check: behaviour said border, now the gradients say border.\n"
        "- In the chart, only the REAL group on the far right has a tall yellow (pad) bar.\n"
        "- Table's last column is the cleanest summary: real ranks first in 7 of 7 pad runs, "
        "0-2 of 7 everywhere else.\n"
        "- Note rescale draws on the most centrally concentrated evidence of the four, 0.373.\n"
        "- Method caveat if asked: seed 0 only.",
    )


def _slide_14_controls(prs: Presentation) -> None:
    """
    Analysis 2: the capacity control and the frequency-domain mechanism.
    """
    slide = _blank(prs, PAPER)
    _chrome(
        slide,
        "Analysis  ·  2 of 2",
        "Not an undertrained CNN — and here is the mechanism",
        title_pt=31,
    )
    _text(
        slide,
        "Is the weak transfer simply a small network trained badly?",
        MARGIN_IN,
        1.76,
        6.10,
        0.32,
        size=13.5,
        color=MUTED,
        bold=True,
    )
    _grid_table(
        slide,
        ["", "Train", "In-domain test", "Cross-generator"],
        [
            ["center_crop", "0.967", "0.908", "0.655"],
            ["rescale", "0.982", "0.902", "0.563"],
        ],
        MARGIN_IN,
        2.18,
        6.10,
        col_ratios=[1.25, 0.85, 1.25, 1.25],
        row_h=0.46,
        header_h=0.46,
        emphasis_col=3,
    )
    _text(
        slide,
        "At 96.7 – 98.2 % train accuracy the network is close to memorizing its 3,600 images, so the "
        "modest in-domain number is a variance limit, not a capacity limit.",
        MARGIN_IN,
        3.72,
        6.10,
        0.74,
        size=13,
        color=SLATE,
        line_spacing=1.24,
    )
    _card(slide, MARGIN_IN, 4.56, 6.10, 1.06, fill=CARD)
    _text(
        slide,
        "rescale fits its training set better, matches center_crop on held-out images of the same "
        "generator, and collapses only when the generator changes.",
        MARGIN_IN + 0.26,
        4.76,
        5.58,
        0.70,
        size=13,
        color=INK,
        bold=True,
        line_spacing=1.24,
    )
    _text(
        slide,
        "Doubling the epoch budget from 40 to 80 on three generators moved in-domain accuracy "
        "+0.020 and cross-generator −0.007. More training buys fit, and nothing in transfer.",
        MARGIN_IN,
        5.76,
        6.10,
        0.72,
        size=12.5,
        color=MUTED,
        line_spacing=1.22,
    )

    _text(
        slide,
        "Radially averaged spectra close the loop",
        7.06,
        1.76,
        5.65,
        0.32,
        size=13.5,
        color=MUTED,
        bold=True,
    )
    _picture(slide, FIGURES_DIR / "spectra.png", 7.06, 2.18, 5.65)
    _text(
        slide,
        "Real minus fake, mean log₁₀ power above 0.35 cycles / pixel",
        7.06,
        3.86,
        5.65,
        0.32,
        size=11.5,
        color=MUTED,
    )
    spectra = [
        ("center_crop", "+0.256", TEAL),
        ("random_crop", "+0.253", TEAL),
        ("rescale", "+0.075", ACCENT),
        ("pad", "+0.058", ACCENT),
    ]
    for index, (name, value, color) in enumerate(spectra):
        left = 7.06 + index * 1.44
        _card(slide, left, 4.26, 1.32, 1.06, fill=CARD)
        _text(slide, value, left, 4.42, 1.32, 0.42, size=19, color=color, font=HEAD_FONT, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, name, left, 4.92, 1.32, 0.28, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    _text(
        slide,
        "Both resampling arms lose roughly two thirds of the high-frequency separation the crops "
        "retain. That is the mechanism behind the transfer collapse.",
        7.06,
        5.52,
        5.65,
        0.80,
        size=13,
        color=INK,
        bold=True,
        line_spacing=1.26,
    )
    _notes(
        slide,
        "1:00  |  Noa  — densest slide, keep moving\n"
        "- Left: kill the obvious objection. Train 0.967 / 0.982 means capacity is not the limit.\n"
        "- The three-level split is the sharpest form of the headline: rescale is FINE in-domain "
        "and only fails cross-generator. A practitioner validating in-distribution sees nothing wrong.\n"
        "- Right: the spectra confirm the assumption we started from. Resampling really does "
        "attenuate the high frequencies that carry the fingerprint.\n"
        "- Four numbers, two groups. Do not narrate the curves.",
    )


def _slide_15_summary(prs: Presentation) -> None:
    """
    Required element: summary and insights.
    """
    slide = _blank(prs, INK)
    _chrome(slide, "Summary & insights", "Three findings that survive our controls", dark=True)
    findings = [
        (
            "A zero-parameter size rule beats every trained detector cross-generator",
            "So cross-generator numbers reported without an explicit size correction are "
            "largely uninterpretable.",
        ),
        (
            "The correction is not a neutral preprocessing detail",
            "Strategies preserving native pixels agree on generator difficulty (ρ = +0.964); "
            "resampling gives an unrelated ordering. A published ranking depends partly on a "
            "choice that is rarely reported.",
        ),
        (
            "A correction can raise accuracy while making the experiment worse",
            "pad trades the size cue for a border cue — found by attribution in all seven of its "
            "runs, and in none elsewhere.",
        ),
    ]
    _bullets(slide, findings, MARGIN_IN, BODY_Y_IN, 8.30, row_h=1.08, gap=0.20, dark=True)
    _card(slide, 9.10, 1.86, 3.62, 1.72, fill=CARD_DARK)
    _text(
        slide,
        "Recommendation",
        9.34,
        2.06,
        3.14,
        0.28,
        size=12,
        color=ACCENT,
        bold=True,
        tracking="60",
    )
    _text(
        slide,
        "Treat size equalization\nas a first-class\nmethodological choice —\nand report more than one.",
        9.34,
        2.40,
        3.14,
        1.00,
        size=13.5,
        color=PAPER,
        line_spacing=1.26,
    )
    _card(slide, 9.10, 3.72, 3.62, 1.66, fill=CARD_DARK)
    _text(
        slide,
        "Limitations",
        9.34,
        3.90,
        3.14,
        0.28,
        size=12,
        color=PAPER_DIM,
        bold=True,
        tracking="60",
    )
    _text(
        slide,
        "One architecture, one input\nresolution. Two seeds. n = 7\ngenerators. Attribution on\nseed 0. ADM unexplained.",
        9.34,
        4.22,
        3.14,
        1.00,
        size=12,
        color=PAPER_DIM,
        line_spacing=1.24,
    )
    _card(slide, MARGIN_IN, 5.90, CONTENT_W_IN, 0.94, fill=CARD_DARK)
    _text(
        slide,
        "Code, run instructions, split definitions and every metrics.json from the 56 runs:",
        MARGIN_IN + 0.32,
        6.08,
        7.0,
        0.32,
        size=13,
        color=PAPER_DIM,
    )
    _text(
        slide,
        REPO_URL,
        MARGIN_IN + 0.32,
        6.40,
        7.0,
        0.32,
        size=14,
        color=PAPER,
        font=MONO_FONT,
        bold=True,
    )
    _text(
        slide,
        "Thank you  —  questions?",
        8.30,
        6.14,
        4.10,
        0.44,
        size=17,
        color=GOLD,
        font=HEAD_FONT,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    _notes(
        slide,
        "1:10  |  Noa\n"
        "- Three findings, one sentence each. Do not re-derive them.\n"
        "- The recommendation is the takeaway you want remembered: report more than one correction.\n"
        "- State limitations plainly, do not apologise for them. ADM being unexplained is honest, "
        "and we say so in the report too.\n"
        "- Repo is public. Invite questions.\n"
        "- Running total to here: 13:05. Backup slides follow if questions need them.",
    )


# ── Backup ────────────────────────────────────────────────────────────────────


def _slide_b1_hyperparameters(prs: Presentation) -> None:
    """
    Backup: the full training configuration and architecture detail.
    """
    slide = _blank(prs, PAPER)
    _chrome(slide, "Backup  ·  B1", "Full training configuration", title_pt=30)
    _grid_table(
        slide,
        ["Setting", "Value", "Setting", "Value"],
        [
            ["Optimizer", "AdamW", "Loss", "BCE-with-logits"],
            ["Learning rate", "3×10⁻⁴", "Precision", "AMP float16"],
            ["Weight decay", "1×10⁻⁴", "Augmentation", "h-flip, p = 0.5"],
            ["Scheduler", "cosine to 0", "Normalization", "(x/255 − 0.5) / 0.5"],
            ["Batch size", "128", "Selection", "best internal val"],
            ["Epochs", "40", "Hardware", "Colab T4, torch 2.11"],
        ],
        MARGIN_IN,
        BODY_Y_IN,
        CONTENT_W_IN,
        col_ratios=[1.1, 1.3, 1.1, 1.4],
        row_h=0.43,
        header_h=0.40,
    )
    _bullets(
        slide,
        [
            (
                "1,173,473 parameters; receptive field 76×76 at the final block, 59 % of the image",
                "Four blocks of [Conv 3×3, BatchNorm, ReLU] ×2 then 2×2 max-pool, widths "
                "32/64/128/256, taking 128×128 down to 8×8; global average pooling; dropout 0.3; "
                "one linear unit.",
            ),
            (
                "Why horizontal flip is the only augmentation",
                "Random resized cropping, JPEG jitter and blur all resample or re-encode the image "
                "— injecting the very artifact under study. Rotations by multiples of 90° are "
                "exact, and were tested separately: off-domain accuracy moved 0.009 over the whole "
                "orbit.",
            ),
        ],
        MARGIN_IN,
        5.14,
        CONTENT_W_IN,
        row_h=0.98,
        gap=0.16,
    )
    _notes(slide, "Backup. Identical across all 56 runs.")


def _slide_b2_adm(prs: Presentation) -> None:
    """
    Backup: the unexplained ADM inversion.
    """
    slide = _blank(prs, PAPER)
    _chrome(slide, "Backup  ·  B2", "The ADM anomaly, reported as unresolved", title_pt=30)
    _stat(slide, "0.863", "ADM in-domain\nunder center_crop", MARGIN_IN, BODY_Y_IN, 3.88, 1.72, value_color=ACCENT)
    _stat(slide, "0.986", "ADM in-domain\nunder rescale", 4.74, BODY_Y_IN, 3.88, 1.72, value_color=TEAL)
    _stat(slide, "0.501", "ADM transfer\nunder rescale", 8.84, BODY_Y_IN, 3.87, 1.72, value_color=ACCENT)
    _bullets(
        slide,
        [
            (
                "Under center_crop ADM is the hardest generator in the study; under rescale it is nearly trivial",
                "And its transfer falls to chance — the lowest figure anywhere in the 7×7.",
            ),
            (
                "The obvious explanation predicts the opposite of what we see",
                "ADM is the only PNG generator in an otherwise JPEG dataset. That asymmetry should "
                "make ADM easiest under center_crop, where native pixels and real-image JPEG "
                "blocking both survive. We observe the reverse.",
            ),
            (
                "ADM's row feeds the ranking, so we report it rather than explain it away",
                "Offering an untested mechanism here would weaken the parts of the analysis we can "
                "support. Re-encoding ADM as JPEG and re-running its row would separate format "
                "from resolution.",
            ),
        ],
        MARGIN_IN,
        4.02,
        CONTENT_W_IN,
        row_h=0.94,
        gap=0.12,
    )
    _notes(slide, "Backup. If asked why we did not explain it: we have not tested a mechanism, so we do not claim one.")


def _slide_b3_baselines(prs: Presentation) -> None:
    """
    Backup: the two baseline transfer matrices in full.
    """
    slide = _blank(prs, PAPER)
    _chrome(slide, "Backup  ·  B3", "The zero-parameter rules, cell by cell", title_pt=30)
    _picture(slide, FIGURES_DIR / "transfer_matrices_baselines_report.png", 2.30, 1.74, 8.70)
    _notes(
        slide,
        "Backup. size_lookup is 1.000 on shared-resolution pairs and 0.500 otherwise, in 47 of 49 "
        "cells; the exceptions are the single 512-square Midjourney row. square_rule is flat 0.977 "
        "everywhere.",
    )


def _slide_b4_spectra(prs: Presentation) -> None:
    """
    Backup: the full radially averaged spectra panel and its method note.
    """
    slide = _blank(prs, PAPER)
    _chrome(slide, "Backup  ·  B4", "Radially averaged spectra, per generator", title_pt=30)
    _picture(slide, FIGURES_DIR / "spectra.png", MARGIN_IN, 1.94, CONTENT_W_IN, max_h=2.96)
    _bullets(
        slide,
        [
            (
                "No window function is applied",
                "The pad arm's zero border is a real feature of that arm, not an artifact of the "
                "analysis, so windowing it away would hide the thing we are measuring.",
            ),
            (
                "Rotations cannot be tested this way",
                "Radially averaged spectra are exactly invariant under multiples of 90°: "
                "(u, v) → (−v, u) preserves radius. The invariance is a theorem, not a finding.",
            ),
        ],
        MARGIN_IN,
        5.24,
        CONTENT_W_IN,
        row_h=0.94,
        gap=0.14,
    )
    _notes(slide, "Backup. Curves are real minus fake, per generator, per strategy.")


# ── Assembly ──────────────────────────────────────────────────────────────────

BUILDERS = (
    _slide_01_title,
    _slide_02_problem,
    _slide_03_confound,
    _slide_04_question,
    _slide_05_related,
    _slide_06_data,
    _slide_07_strategies,
    _slide_08_model,
    _slide_09_baselines,
    _slide_10_matrices,
    _slide_11_headline,
    _slide_12_pad,
    _slide_13_attribution,
    _slide_14_controls,
    _slide_15_summary,
    _slide_b1_hyperparameters,
    _slide_b2_adm,
    _slide_b3_baselines,
    _slide_b4_spectra,
)


def build_deck(out_path: Path) -> Path:
    """
    Assemble every slide and write the presentation.

    Args:
        out_path: Destination .pptx path.

    Returns:
        The path written.
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    for builder in BUILDERS:
        builder(prs)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main() -> None:
    """
    Build the deck and report where it landed.
    """
    path = build_deck(OUTPUT_PATH)
    print(f"wrote {path}  ({len(BUILDERS)} slides: 15 content + 4 backup)")


if __name__ == "__main__":
    main()
